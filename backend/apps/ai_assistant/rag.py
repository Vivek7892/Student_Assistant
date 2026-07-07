import os
import json
import logging
import tiktoken
from pathlib import Path
from django.conf import settings
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import ConversationalRetrievalChain
from langchain.prompts import ChatPromptTemplate, SystemMessagePromptTemplate, HumanMessagePromptTemplate
from langchain.schema import HumanMessage
import pypdf
import docx
from pptx import Presentation

logger = logging.getLogger(__name__)

ENCODING = tiktoken.get_encoding('cl100k_base')


def count_tokens(text: str) -> int:
    return len(ENCODING.encode(text))


def truncate_to_tokens(text: str, max_tokens: int) -> str:
    tokens = ENCODING.encode(text)
    return ENCODING.decode(tokens[:max_tokens]) if len(tokens) > max_tokens else text


class DocumentProcessor:
    SUPPORTED_TYPES = {'.pdf', '.docx', '.pptx', '.txt'}

    def extract_text(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext == '.docx':
            return self._extract_docx(file_path)
        elif ext == '.pptx':
            return self._extract_pptx(file_path)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        raise ValueError(f'Unsupported file type: {ext}')

    def _extract_pdf(self, file_path: str) -> str:
        text = []
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)

    def _extract_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)

    def _extract_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return '\n'.join(text)


class RAGPipeline:
    # Token budget: gpt-4o context = 128k, keep responses focused
    CHUNK_TOKENS = 512        # ~512 tokens per chunk (tighter = more precise retrieval)
    CHUNK_OVERLAP_TOKENS = 64
    MAX_CONTEXT_TOKENS = 6000  # max tokens sent as context to LLM
    TOP_K = 6                  # retrieve top-6 chunks, then MMR re-ranks

    def __init__(self):
        self._llm = None
        self._embeddings = None
        # Token-aware splitter: approximate chars from tokens (1 token ≈ 4 chars)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.CHUNK_TOKENS * 4,
            chunk_overlap=self.CHUNK_OVERLAP_TOKENS * 4,
            length_function=count_tokens,  # use tiktoken for accurate splitting
            separators=['\n\n', '\n', '. ', ' ', ''],
        )
        self.persist_dir = str(settings.CHROMA_PERSIST_DIR)
        os.makedirs(self.persist_dir, exist_ok=True)

    @property
    def llm(self):
        if not self._llm:
            key = settings.OPENAI_API_KEY
            if not key:
                raise ValueError('OpenAI API key not configured')
            self._llm = ChatOpenAI(
                model='gpt-4o',
                temperature=0.2,
                max_tokens=1024,
                openai_api_key=key,
                streaming=False,
            )
        return self._llm

    @property
    def embeddings(self):
        if not self._embeddings:
            key = settings.OPENAI_API_KEY
            if not key:
                raise ValueError('OpenAI API key not configured')
            self._embeddings = OpenAIEmbeddings(
                model='text-embedding-3-small',  # faster + cheaper than ada-002
                openai_api_key=key,
            )
        return self._embeddings

    def get_vectorstore(self, collection_name: str):
        return Chroma(
            collection_name=collection_name,
            embedding_function=self.embeddings,
            persist_directory=self.persist_dir,
        )

    def index_document(self, text: str, collection_name: str, metadata: dict = None) -> int:
        chunks = self.text_splitter.split_text(text)
        metadatas = [metadata or {} for _ in chunks]
        vectorstore = self.get_vectorstore(collection_name)
        vectorstore.add_texts(texts=chunks, metadatas=metadatas)
        vectorstore.persist()
        logger.info(f'Indexed {len(chunks)} chunks into {collection_name}')
        return len(chunks)

    def query(self, question: str, collection_name: str, chat_history: list = None, language: str = 'english') -> dict:
        vectorstore = self.get_vectorstore(collection_name)

        # MMR retrieval: balances relevance + diversity
        retriever = vectorstore.as_retriever(
            search_type='mmr',
            search_kwargs={'k': self.TOP_K, 'fetch_k': self.TOP_K * 3, 'lambda_mult': 0.7},
        )

        lang_note = 'Respond in Kannada.' if language == 'kannada' else ''

        system_template = (
            'You are an expert AI study tutor. Answer the student\'s question using ONLY the provided context. '
            'Be precise, structured, and educational. If the answer is not in the context, say so clearly. '
            'Format your response with clear sections when helpful. Cite source material when relevant. '
            + lang_note +
            '\n\nContext:\n{context}'
        )

        messages = [
            SystemMessagePromptTemplate.from_template(system_template),
            HumanMessagePromptTemplate.from_template('{question}'),
        ]
        qa_prompt = ChatPromptTemplate.from_messages(messages)

        chain = ConversationalRetrievalChain.from_llm(
            llm=self.llm,
            retriever=retriever,
            return_source_documents=True,
            combine_docs_chain_kwargs={'prompt': qa_prompt},
            verbose=False,
        )

        history = chat_history or []
        result = chain({'question': question, 'chat_history': history})

        sources = []
        seen = set()
        for doc in result.get('source_documents', []):
            snippet = doc.page_content[:200]
            if snippet not in seen:
                seen.add(snippet)
                sources.append({'content': snippet, 'metadata': doc.metadata})

        tokens_used = count_tokens(question) + count_tokens(result['answer'])

        return {
            'answer': result['answer'],
            'sources': sources,
            'tokens_used': tokens_used,
        }

    def generate_quiz(self, text: str, num_questions: int = 10, difficulty: str = 'medium') -> list:
        # Truncate to token budget before sending
        text = truncate_to_tokens(text, 5000)
        prompt = (
            f'Generate {num_questions} multiple-choice questions at {difficulty} difficulty '
            f'from the study material below.\n\n'
            f'Return ONLY a valid JSON array:\n'
            f'[{{"question":"...","options":["A) ...","B) ...","C) ...","D) ..."],'
            f'"correct_answer":"A","explanation":"..."}}]\n\n'
            f'Study Material:\n{text}'
        )
        response = self.llm.invoke([HumanMessage(content=prompt)])
        return self._parse_json(response.content, fallback=[])

    def generate_flashcards(self, text: str, num_cards: int = 15) -> list:
        text = truncate_to_tokens(text, 5000)
        prompt = (
            f'Generate {num_cards} flashcards from the study material below.\n\n'
            f'Return ONLY a valid JSON array:\n'
            f'[{{"front":"Key term or question","back":"Definition or answer"}}]\n\n'
            f'Study Material:\n{text}'
        )
        response = self.llm.invoke([HumanMessage(content=prompt)])
        return self._parse_json(response.content, fallback=[])

    def summarize(self, text: str) -> str:
        text = truncate_to_tokens(text, 8000)
        prompt = (
            'Summarize the following study material with:\n'
            '1. Brief Overview (2-3 sentences)\n'
            '2. Key Concepts (bullet points)\n'
            '3. Important Details\n'
            '4. Exam Tips\n\n'
            f'Material:\n{text}'
        )
        response = self.llm.invoke([HumanMessage(content=prompt)])
        return response.content

    def generate_study_plan(self, subjects: list, duration_days: int, exam_date: str) -> dict:
        subjects_text = '\n'.join([f'- {s}' for s in subjects])
        prompt = (
            f'Create a {duration_days}-day study plan for exam on {exam_date}.\n\n'
            f'Subjects:\n{subjects_text}\n\n'
            f'Return JSON:\n'
            f'{{"overview":"...","daily_schedule":[{{"day":1,"date":"YYYY-MM-DD",'
            f'"tasks":[{{"subject":"...","topic":"...","duration_hours":2,"type":"study"}}]}}],'
            f'"tips":["tip1"]}}'
        )
        response = self.llm.invoke([HumanMessage(content=prompt)])
        return self._parse_json(response.content, fallback={'overview': response.content, 'daily_schedule': [], 'tips': []})

    def _parse_json(self, content: str, fallback):
        try:
            content = content.strip()
            if content.startswith('```'):
                parts = content.split('```')
                content = parts[1]
                if content.startswith('json'):
                    content = content[4:]
            return json.loads(content.strip())
        except json.JSONDecodeError:
            logger.error('Failed to parse LLM JSON response')
            return fallback


rag_pipeline = RAGPipeline()
doc_processor = DocumentProcessor()
