"""
Gemini RAG Pipeline
- LLM       : Gemini 1.5 Flash (fast) / Pro (quality)
- Embeddings: text-embedding-004
- Vector DB : ChromaDB (persisted, per-material collections)
- Chunking  : tiktoken-aware RecursiveCharacterTextSplitter
- Retrieval : MMR top-6, falls back to general Gemini knowledge if no docs
- Streaming : generator-based for WebSocket token streaming
"""

import os
import json
import time
import logging
import tiktoken
from pathlib import Path
from typing import Generator

from google import genai
from google.genai import types
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.embeddings import Embeddings
from django.conf import settings

import pypdf
import docx
from pptx import Presentation

logger = logging.getLogger(__name__)

_ENCODER = tiktoken.get_encoding('cl100k_base')


def count_tokens(text: str) -> int:
    return len(_ENCODER.encode(text))


def truncate_to_tokens(text: str, max_tokens: int) -> str:
    tokens = _ENCODER.encode(text)
    return _ENCODER.decode(tokens[:max_tokens]) if len(tokens) > max_tokens else text


# ── Document Processor ────────────────────────────────────────────────────────

class DocumentProcessor:
    SUPPORTED_TYPES = {'.pdf', '.docx', '.pptx', '.txt'}

    def extract_text(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        extractors = {
            '.pdf':  self._extract_pdf,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.txt':  self._extract_txt,
        }
        if ext not in extractors:
            raise ValueError(f'Unsupported file type: {ext}')
        return extractors[ext](file_path)

    def _extract_pdf(self, path: str) -> str:
        with open(path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            return '\n'.join(page.extract_text() or '' for page in reader.pages)

    def _extract_docx(self, path: str) -> str:
        doc = docx.Document(path)
        return '\n'.join(p.text for p in doc.paragraphs)

    def _extract_pptx(self, path: str) -> str:
        prs = Presentation(path)
        return '\n'.join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, 'text')
        )

    def _extract_txt(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


# ── Native Gemini Embeddings (replaces langchain-google-genai) ───────────────

class GenAIEmbeddings(Embeddings):
    """LangChain-compatible embeddings using the new google-genai SDK."""

    def __init__(self, api_key: str):
        self._client = genai.Client(api_key=api_key)

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        return [
            self._client.models.embed_content(
                model='models/text-embedding-004',
                contents=t,
                config=types.EmbedContentConfig(task_type='RETRIEVAL_DOCUMENT'),
            ).embeddings[0].values
            for t in texts
        ]

    def embed_query(self, text: str) -> list[float]:
        return self._client.models.embed_content(
            model='models/text-embedding-004',
            contents=text,
            config=types.EmbedContentConfig(task_type='RETRIEVAL_QUERY'),
        ).embeddings[0].values


# ── Gemini RAG Pipeline ───────────────────────────────────────────────────────

class GeminiRAGPipeline:
    CHUNK_TOKENS         = 512
    CHUNK_OVERLAP_TOKENS = 64
    TOP_K                = 6
    MAX_CONTEXT_TOKENS   = 8000

    SAFETY_SETTINGS = [
        types.SafetySetting(category='HARM_CATEGORY_HARASSMENT',        threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_HATE_SPEECH',       threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_SEXUALLY_EXPLICIT', threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_DANGEROUS_CONTENT', threshold='BLOCK_NONE'),
    ]

    def __init__(self):
        self._client     = None
        self._embeddings = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.CHUNK_TOKENS * 4,
            chunk_overlap=self.CHUNK_OVERLAP_TOKENS * 4,
            length_function=count_tokens,
            separators=['\n\n', '\n', '. ', ' ', ''],
        )
        self.persist_dir = str(settings.CHROMA_PERSIST_DIR)
        os.makedirs(self.persist_dir, exist_ok=True)

    def _get_api_key(self) -> str:
        key = getattr(settings, 'GEMINI_API_KEY', '') or os.environ.get('GEMINI_API_KEY', '')
        if not key:
            raise ValueError('GEMINI_API_KEY not set. Get a key from https://aistudio.google.com/app/apikey and add it to backend/.env')
        if not key.startswith('AIza'):
            raise ValueError(
                f'GEMINI_API_KEY looks invalid (starts with "{key[:6]}..."). '
                'It must start with "AIza". Get a valid key from https://aistudio.google.com/app/apikey'
            )
        return key

    @property
    def client(self) -> genai.Client:
        if not self._client:
            self._client = genai.Client(api_key=self._get_api_key())
        return self._client

    @property
    def _gen_config(self) -> types.GenerateContentConfig:
        return types.GenerateContentConfig(
            temperature=0.4,
            max_output_tokens=4096,
            top_p=0.95,
            top_k=40,
            safety_settings=self.SAFETY_SETTINGS,
        )

    @property
    def _json_config(self) -> types.GenerateContentConfig:
        return types.GenerateContentConfig(
            temperature=0.2,
            max_output_tokens=4096,
            response_mime_type='application/json',
            safety_settings=self.SAFETY_SETTINGS,
        )

    @property
    def embeddings(self) -> 'GenAIEmbeddings':
        if not self._embeddings:
            self._embeddings = GenAIEmbeddings(api_key=self._get_api_key())
        return self._embeddings

    def get_vectorstore(self, collection_name: str) -> Chroma:
        return Chroma(
            collection_name=collection_name,
            embedding_function=self.embeddings,
            persist_directory=self.persist_dir,
        )

    # ── Indexing ──────────────────────────────────────────────────────────────

    def index_document(self, text: str, collection_name: str, metadata: dict = None) -> int:
        chunks = self.text_splitter.split_text(text)
        metadatas = [metadata or {} for _ in chunks]
        vs = self.get_vectorstore(collection_name)
        vs.add_texts(texts=chunks, metadatas=metadatas)
        vs.persist()
        logger.info(f'[RAG] Indexed {len(chunks)} chunks → {collection_name}')
        return len(chunks)

    # ── Retrieval helpers ─────────────────────────────────────────────────────

    def _retrieve_context(self, question: str, collection_name: str) -> tuple[str, list[dict]]:
        """Returns (context_text, sources_list). Returns ('', []) if collection empty."""
        try:
            vs = self.get_vectorstore(collection_name)
            # Check if collection has any documents
            count = vs._collection.count()
            if count == 0:
                return '', []

            retriever = vs.as_retriever(
                search_type='mmr',
                search_kwargs={
                    'k': self.TOP_K,
                    'fetch_k': self.TOP_K * 3,
                    'lambda_mult': 0.7,
                },
            )
            docs = retriever.get_relevant_documents(question)

            context_parts, sources, total_tokens, seen = [], [], 0, set()
            for doc in docs:
                snippet = doc.page_content
                tokens = count_tokens(snippet)
                if total_tokens + tokens > self.MAX_CONTEXT_TOKENS:
                    break
                context_parts.append(snippet)
                total_tokens += tokens
                key = snippet[:100]
                if key not in seen:
                    seen.add(key)
                    sources.append({'content': snippet[:300], 'metadata': doc.metadata})

            return '\n\n---\n\n'.join(context_parts), sources
        except Exception as e:
            logger.warning(f'[RAG] Context retrieval failed for {collection_name}: {e}')
            return '', []

    def _build_prompt(self, question: str, context: str, history_text: str, language: str) -> str:
        lang_note = 'Respond in Kannada.' if language == 'kannada' else 'Respond in English.'

        if context:
            context_section = f"""
UPLOADED STUDY MATERIAL CONTEXT (use this as primary source):
{context}

---
"""
            source_instruction = (
                "- Prioritize answers from the uploaded study material context above.\n"
                "- If the question is answered by the context, cite it.\n"
                "- If the context doesn't fully answer the question, supplement with your general knowledge and clearly indicate which parts come from general knowledge vs the uploaded material.\n"
            )
        else:
            context_section = ""
            source_instruction = (
                "- No study materials have been uploaded yet. Answer from your general knowledge as an expert tutor.\n"
                "- Mention that the student can upload study materials (PDF, DOCX, PPTX) for more personalized answers.\n"
            )

        return f"""You are an expert AI study tutor for university students, powered by Google Gemini.

INSTRUCTIONS:
- Be helpful, accurate, and educational.
- Use markdown formatting: **bold** for key terms, bullet points for lists, numbered steps for processes, code blocks for code.
- Give thorough, well-structured answers.
- {lang_note}
{source_instruction}
{context_section}
CONVERSATION HISTORY:
{history_text or 'No previous messages.'}

STUDENT QUESTION: {question}

ANSWER:"""

    def _call_with_retry(self, fn, *args, max_retries: int = 3, **kwargs):
        """Call fn(*args, **kwargs) with exponential backoff on 429 rate-limit errors."""
        for attempt in range(max_retries):
            try:
                return fn(*args, **kwargs)
            except Exception as e:
                err = str(e)
                is_rate_limit = '429' in err or 'quota' in err.lower() or 'rate' in err.lower()
                if is_rate_limit and attempt < max_retries - 1:
                    wait = 2 ** (attempt + 1)  # 2s, 4s, 8s
                    logger.warning(f'[RAG] Rate limited, retrying in {wait}s (attempt {attempt + 1}/{max_retries})')
                    time.sleep(wait)
                else:
                    raise

    # ── Query (non-streaming) ─────────────────────────────────────────────────

    def query(self, question: str, collection_name: str,
              chat_history: list = None, language: str = 'english') -> dict:
        context, sources = self._retrieve_context(question, collection_name)

        history_text = '\n'.join(
            f'Student: {u}\nTutor: {a}' for u, a in (chat_history or [])
        )

        prompt = self._build_prompt(question, context, history_text, language)
        response = self._call_with_retry(
            self.client.models.generate_content,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._gen_config
        )

        tokens_used = count_tokens(prompt) + count_tokens(response.text)
        return {
            'answer': response.text,
            'sources': sources,
            'tokens_used': tokens_used,
            'model': 'gemini-2.0-flash-lite',
        }

    # ── Query (streaming) ─────────────────────────────────────────────────────

    def query_stream(self, question: str, collection_name: str,
                     chat_history: list = None, language: str = 'english') -> Generator[dict, None, None]:
        context, sources = self._retrieve_context(question, collection_name)

        history_text = '\n'.join(
            f'Student: {u}\nTutor: {a}' for u, a in (chat_history or [])
        )

        prompt = self._build_prompt(question, context, history_text, language)

        yield {'type': 'sources', 'data': sources}

        full_answer = ''
        for chunk in self._call_with_retry(
            self.client.models.generate_content_stream,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._gen_config
        ):
            if chunk.text:
                full_answer += chunk.text
                yield {'type': 'token', 'data': chunk.text}

        tokens_used = count_tokens(prompt) + count_tokens(full_answer)
        yield {'type': 'done', 'tokens_used': tokens_used, 'model': 'gemini-2.0-flash-lite'}

    # ── Quiz generation ───────────────────────────────────────────────────────

    def generate_quiz(self, text: str, num_questions: int = 10, difficulty: str = 'medium') -> list:
        text = truncate_to_tokens(text, 6000)
        prompt = (
            f'Generate {num_questions} multiple-choice questions at {difficulty} difficulty '
            f'from the study material below.\n\n'
            f'Return ONLY valid JSON array (no markdown):\n'
            f'[{{"question":"...","options":["A) ...","B) ...","C) ...","D) ..."],'
            f'"correct_answer":"A","explanation":"..."}}]\n\n'
            f'Study Material:\n{text}'
        )
        response = self._call_with_retry(
            self.client.models.generate_content,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._json_config
        )
        return {
            'items': self._parse_json(response.text, fallback=[]),
            'tokens_used': count_tokens(prompt) + count_tokens(response.text or ''),
            'model': 'gemini-2.0-flash-lite',
        }

    # ── Flashcard generation ──────────────────────────────────────────────────

    def generate_flashcards(self, text: str, num_cards: int = 15) -> list:
        text = truncate_to_tokens(text, 6000)
        prompt = (
            f'Generate {num_cards} flashcards from the study material below.\n\n'
            f'Return ONLY valid JSON array (no markdown):\n'
            f'[{{"front":"Key term or question","back":"Definition or answer"}}]\n\n'
            f'Study Material:\n{text}'
        )
        response = self._call_with_retry(
            self.client.models.generate_content,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._json_config
        )
        return {
            'items': self._parse_json(response.text, fallback=[]),
            'tokens_used': count_tokens(prompt) + count_tokens(response.text or ''),
            'model': 'gemini-2.0-flash-lite',
        }

    # ── Summarize ─────────────────────────────────────────────────────────────

    def summarize(self, text: str) -> str:
        text = truncate_to_tokens(text, 10000)
        prompt = (
            'Summarize the following study material with:\n'
            '1. **Overview** (2-3 sentences)\n'
            '2. **Key Concepts** (bullet points)\n'
            '3. **Important Details**\n'
            '4. **Exam Tips**\n\n'
            f'Material:\n{text}'
        )
        response = self._call_with_retry(
            self.client.models.generate_content,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._gen_config
        )
        return response.text

    # ── Study plan ────────────────────────────────────────────────────────────

    def generate_study_plan(self, subjects: list, duration_days: int, exam_date: str) -> dict:
        subjects_text = '\n'.join(f'- {s}' for s in subjects)
        prompt = (
            f'Create a {duration_days}-day study plan for exam on {exam_date}.\n\n'
            f'Subjects:\n{subjects_text}\n\n'
            f'Return ONLY valid JSON (no markdown):\n'
            f'{{"overview":"...","daily_schedule":[{{"day":1,"date":"YYYY-MM-DD",'
            f'"tasks":[{{"subject":"...","topic":"...","duration_hours":2,"type":"study"}}]}}],'
            f'"tips":["tip1"]}}'
        )
        response = self._call_with_retry(
            self.client.models.generate_content,
            model='gemini-2.0-flash-lite', contents=prompt, config=self._gen_config
        )
        return self._parse_json(
            response.text,
            fallback={'overview': response.text, 'daily_schedule': [], 'tips': []}
        )

    # ── JSON helper ───────────────────────────────────────────────────────────

    def _parse_json(self, content: str, fallback):
        try:
            content = content.strip()
            if '```' in content:
                parts = content.split('```')
                content = parts[1] if len(parts) > 1 else content
                if content.startswith('json'):
                    content = content[4:]
            return json.loads(content.strip())
        except (json.JSONDecodeError, IndexError):
            logger.error('[RAG] Failed to parse JSON response')
            return fallback


# ── Singletons ────────────────────────────────────────────────────────────────
gemini_rag    = GeminiRAGPipeline()
doc_processor = DocumentProcessor()
